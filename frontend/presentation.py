import os
import requests
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import comtypes.client
import pythoncom
from PIL import Image
import tempfile
from typing import List, Dict

# Добавляем текстовое поле для содержимого слайда
def add_textbox_with_text(slide, text, left, top, width, height):
    left = Inches(left)
    top = Inches(top)
    width = Inches(width)
    height = Inches(height)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = text

    # Перенос текста в рамках текстового блока
    text_frame.fit_text()

    #text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT


def search_unsplash_images(query: str):
    """
    Выполняет поиск изображений на Unsplash по заданному запросу и возвращает объект ответа.

    Args:
        query (str): Поисковый запрос.

    Returns:
        response (requests.Response): Ответ от Unsplash API.
    """
    ACCESS_KEY = 'ns_k09as7WL-3HKDhcgCLu7fDNC0eZ-53klLpMg3eDk'
    url = 'https://api.unsplash.com/search/photos'
    params = {
        'query': query,
        'client_id': ACCESS_KEY,
        'lang': 'ru'
    }
    response = requests.get(url, params=params)
    # print(response.json())
    return response


def save_image_from_url(url: str, save_path: str) -> None:
    """
    Сохраняет изображение по указанному URL на диск по заданному пути.

    Args:
        url (str): Ссылка на изображение.
        save_path (str): Путь, по которому сохранить изображение.

    Пример:
        save_image_from_url('https://example.com/image.jpg', 'local_image.jpg')
    """
    response = requests.get(url)
    if response.status_code == 200:
        with open(save_path, 'wb') as f:
            f.write(response.content)
        print(f"Изображение успешно сохранено по пути: {save_path}")
    else:
        print(f"Не удалось скачать изображение. Код ответа: {response.status_code}")


def create_presentation(topic: str, slides_data: List[Dict], background_path: str) -> str:
    """
    Создаёт презентацию PPTX с заданной темой, слайдами и фоном.

    Аргументы:
        topic (str): Название презентации (будет именем файла).
        slides_data (List[Dict]): Список словарей с ключами "title" и "content" для каждого слайда.
        background_path (str): Путь к изображению фона для слайдов.

    Возвращает:
        str: Путь к сохранённому файлу презентации PPTX.
    """
    # Создаем папку для изображений
    os.makedirs("temp_images", exist_ok=True)

    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[1])

    # Устанавливаем фоновое изображение на весь слайд
    title_slide.shapes.add_picture(
        background_path,
        left=0,
        top=0,
        width=prs.slide_width,
        height=prs.slide_height
    )

    # Устанавливаем заголовок слайда
    add_textbox_with_text(title_slide, slides_data["Presentation_title"], 3.5, 3, 4, 2)

    for slide_id, slide_data in enumerate(slides_data["Slides"]):
        # Добавляем слайд с макетом "Заголовок и объект"
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Устанавливаем фоновое изображение на весь слайд
        slide.shapes.add_picture(
            background_path,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )

        # Устанавливаем заголовок слайда
        add_textbox_with_text(slide, slide_data["Slide_title"], 3, 0.8, 6, 2)

        # Устанавливаем текст слайда
        add_textbox_with_text(slide, slide_data["Slide_content"], 1, 3, 4, 4)

        # Поиск и добавление изображения
        text_search = slide_data["Slide_title"] + slide_data["Slide_content"]

        image_url = search_unsplash_images(text_search).json()['results'][0]['urls']['raw']

        image_path = f"temp_images/slide_{slide_id+1}.jpg"

        save_image_from_url(image_url, image_path)
        
        left_img = Inches(5.5)
        top_img = Inches(3)
        width_img = Inches(4)
        height_img = Inches(4)
        
        slide.shapes.add_picture(
            image_path,
            left_img, top_img,
            width_img, height_img
        )

    pptx_path = f"{topic}.pptx"
    prs.save(pptx_path)

    # Удаляем временные изображения
    for file in os.listdir("temp_images"):
        os.remove(os.path.join("temp_images", file))
    os.rmdir("temp_images")

    return pptx_path


def convert_to_pdf(pptx_path: str) -> str:
    """
    Конвертирует презентацию PPTX в PDF.

    Args:
        pptx_path (str): Путь к файлу презентации PPTX

    Returns:
        str: Путь к созданному PDF файлу

    Raises:
        Exception: Если произошла ошибка при конвертации
    """
    pdf_path = pptx_path.replace(".pptx", ".pdf")

    try:
        pythoncom.CoInitialize()
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        abs_pptx_path = os.path.abspath(pptx_path)
        abs_pdf_path = os.path.abspath(pdf_path)
        deck = powerpoint.Presentations.Open(abs_pptx_path)
        deck.SaveAs(abs_pdf_path, 32)  # 32 - формат PDF
        deck.Close()
        powerpoint.Quit()
        return pdf_path
    
    except Exception as e:
        raise Exception(f"Ошибка конвертации: {str(e)}")
    finally:
        pythoncom.CoUninitialize()


def generate_previews(pptx_path: str, slides_count: int) -> List[Image.Image]:
    """
    Генерирует превью (изображения) слайдов презентации PPTX.

    Args:
        pptx_path (str): Путь к файлу презентации PPTX.
        slides_count (int): Количество слайдов в презентации.

    Returns:
        List[Image.Image]: Список изображений-превью для каждого слайда.

    Raises:
        Exception: Если произошла ошибка при генерации превью.
    """
    previews = []
    try:
        pythoncom.CoInitialize()
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        abs_pptx_path = os.path.abspath(pptx_path)
        deck = powerpoint.Presentations.Open(abs_pptx_path)

        temp_dir = tempfile.mkdtemp()
        for i in range(1, slides_count + 1):
            img_filename = f"slide_{i}.jpg"
            img_path = os.path.join(temp_dir, img_filename)
            deck.Slides(i).Export(img_path, "JPG")
            previews.append(Image.open(img_path))

        deck.Close()
        powerpoint.Quit()
        return previews

    except Exception as e:
        raise Exception(f"Ошибка генерации превью: {str(e)}")
    finally:
        pythoncom.CoUninitialize()