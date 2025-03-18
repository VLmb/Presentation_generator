import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from io import BytesIO
import os
import tempfile
from pdf2image import convert_from_path
from PIL import Image


# Функция для создания презентации
def create_presentation(topic, slides_data, background_image):
    prs = Presentation()

    # Устанавка фона для всех слайдов
    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # пустой макет
        if background_image:
            # Добавление изображения в качестве фона
            slide.shapes.add_picture(background_image, 0, 0, width=prs.slide_width, height=prs.slide_height)

        # Добавление заголовка и текста
        title = slide.shapes.title
        title.text = slide_data["title"]
        content = slide.shapes.placeholders[1].text_frame
        content.text = slide_data["content"]

        # Добавление изображения, если оно есть
        if slide_data["image"]:
            img_path = slide_data["image"]
            slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(4))

    # Сохраняем презентацию
    prs.save(f"{topic}_presentation.pptx")
    return f"{topic}_presentation.pptx"


# конвертация PPTX в PDF
def convert_pptx_to_pdf(pptx_path):
    pdf_path = pptx_path.replace(".pptx", ".pdf")
    os.system(f"libreoffice --headless --convert-to pdf {pptx_path} --outdir {os.path.dirname(pptx_path)}")
    return pdf_path


# Создание миниатюр слайдов
def create_slide_previews(pptx_path):
    images = convert_from_path(pptx_path)
    return images


# Основной интерфейс
st.title("Генератор презентаций")
st.write("Введите данные для создания презентации:")

# Ввод данных пользователем
topic = st.text_input("Тема презентации")
num_slides = st.number_input("Количество слайдов", min_value=1, max_value=20, value=3)

# Выбор фона
background_options = {
    "Фон 1": "background1.jpg",
    "Фон 2": "background2.jpg",
    "Фон 3": "background3.jpg"
}
selected_background = st.selectbox("Выберите фон", list(background_options.keys()))
background_image = background_options[selected_background]

slides_data = []

# Ввод данных
for i in range(num_slides):
    st.subheader(f"Слайд {i + 1}")
    title = st.text_input(f"Заголовок слайда {i + 1}", key=f"title_{i}")
    content = st.text_area(f"Текст слайда {i + 1}", key=f"content_{i}")
    image = st.file_uploader(f"Загрузите изображение для слайда {i + 1}", type=["png", "jpg", "jpeg"], key=f"image_{i}")
    slides_data.append({"title": title, "content": content, "image": image})

if st.button("Сгенерировать презентацию"):
    if topic and slides_data:
        # Сохранение загруженных изображений во временные файлы
        for slide_data in slides_data:
            if slide_data["image"]:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_file:
                    tmp_file.write(slide_data["image"].getvalue())
                    slide_data["image"] = tmp_file.name

        presentation_path = create_presentation(topic, slides_data, background_image)
        st.success("Презентация успешно создана!")

        # Миниатюры слайдов
        slide_previews = create_slide_previews(presentation_path)
        st.subheader("Предпросмотр слайдов:")
        for i, image in enumerate(slide_previews):
            st.image(image, caption=f"Слайд {i + 1}", use_column_width=True)

        # Отображение ссылки для скачивания PPTX
        with open(presentation_path, "rb") as file:
            btn_pptx = st.download_button(
                label="Скачать презентацию (PPTX)",
                data=file,
                file_name=presentation_path,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        # Конвертация в PDF
        pdf_path = convert_pptx_to_pdf(presentation_path)
        with open(pdf_path, "rb") as file:
            btn_pdf = st.download_button(
                label="Скачать презентацию (PDF)",
                data=file,
                file_name=pdf_path,
                mime="application/pdf"
            )
    else:
        st.error("Пожалуйста, заполните все поля.")