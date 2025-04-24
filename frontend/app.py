import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import tempfile
import os
import comtypes.client
import pythoncom
from PIL import Image

BACKGROUNDS = {
    "Синий градиент": "backgrounds/blue_gradient.jpg",
    "Тёмный узор": "backgrounds/dark_pattern.jpg",
    "Светлый мрамор": "backgrounds/light_marble.jpg"
}


def init_environment():
    os.makedirs("backgrounds", exist_ok=True)
    return True


def create_presentation(topic, slides_data, background_path):
    prs = Presentation()
    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.add_picture(background_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        title = slide.shapes.title
        title.text = slide_data["title"]
        left = Inches(1)
        top = Inches(1.5)
        width = prs.slide_width - Inches(2)
        height = prs.slide_height - Inches(2)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = slide_data["content"]
    pptx_path = f"{topic}.pptx"
    prs.save(pptx_path)
    return pptx_path


def convert_to_pdf(pptx_path):
    pdf_path = pptx_path.replace(".pptx", ".pdf")
    try:
        pythoncom.CoInitialize()
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
        deck.SaveAs(os.path.abspath(pdf_path), 32)  # 32 - формат PDF
        deck.Close()
        powerpoint.Quit()
        return pdf_path
    except Exception as e:
        st.error(f"Ошибка конвертации: {str(e)}")
        return None
    finally:
        pythoncom.CoUninitialize()


def generate_previews(pptx_path, slides_count):
    previews = []
    try:
        pythoncom.CoInitialize()
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path))

        temp_dir = tempfile.mkdtemp()
        for i in range(1, slides_count + 1):
            img_path = os.path.join(temp_dir, f"slide_{i}.jpg")
            deck.Slides(i).Export(img_path, "JPG")
            previews.append(Image.open(img_path))

        deck.Close()
        powerpoint.Quit()
        return previews
    except Exception as e:
        st.error(f"Ошибка генерации превью: {str(e)}")
        return []
    finally:
        pythoncom.CoUninitialize()


def parse_file_to_slides(uploaded_file):
    content = uploaded_file.getvalue().decode("utf-8")
    slides = content.split("\n\n")
    return [{"title": f"Слайд {i + 1}", "content": slide} for i, slide in enumerate(slides)]


init_environment()

st.title("📊 Генератор презентаций")
mode = st.radio("Режим работы:", ["Обычный генератор", "Генератор по файлу"])

col1, col2 = st.columns(2)
with col1:
    topic = st.text_input("Название презентации", "Моя презентация")
with col2:
    selected_bg = st.selectbox("Фон презентации", list(BACKGROUNDS.keys()))
bg_path = BACKGROUNDS[selected_bg]

if mode == "Обычный генератор":
    slides_count = st.number_input("Количество слайдов", 1, 50, 3)
    slides_data = [{"title": f"Слайд {i + 1}", "content": ""} for i in range(slides_count)]
    for i, slide in enumerate(slides_data):
        with st.expander(f"Слайд {i + 1}"):
            slide["title"] = st.text_input("Заголовок", slide["title"], key=f"title_{i}")
            slide["content"] = st.text_area("Содержание", slide["content"], key=f"content_{i}")
else:
    uploaded_file = st.file_uploader("Загрузите текстовый файл", type=["txt"])
    if uploaded_file:
        slides_data = parse_file_to_slides(uploaded_file)
        st.success(f"Загружено {len(slides_data)} слайдов")
    else:
        slides_data = []

if st.button("Создать презентацию", type="primary") and slides_data:
    with st.spinner("Идёт создание презентации..."):
        pptx_path = create_presentation(topic, slides_data, bg_path)
        pdf_path = convert_to_pdf(pptx_path)
        previews = generate_previews(pptx_path, len(slides_data))

    st.success("Готово!")

    if previews:
        st.subheader("Превью слайдов")
        cols = st.columns(3)
        for i, img in enumerate(previews):
            cols[i % 3].image(img, caption=f"Слайд {i + 1}", use_column_width=True)

    col1, col2 = st.columns(2)
    with col1:
        with open(pptx_path, "rb") as f:
            st.download_button("📥 Скачать PPTX", f.read(), file_name=os.path.basename(pptx_path),
                               mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    with col2:
        if pdf_path and os.path.exists(pdf_path):
            with open(pdf_path, "rb") as f:
                st.download_button("📥 Скачать PDF", f.read(), file_name=os.path.basename(pdf_path),
                                   mime="application/pdf")