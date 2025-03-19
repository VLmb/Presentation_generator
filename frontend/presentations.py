import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from io import BytesIO
import os
import tempfile
from pdf2image import convert_from_path
from PIL import Image


# Функция для создания презентации
def create_presentation(topic, slides_data, background_image):
    prs = Presentation()

    # Устанавка фона для всех слайдов
    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # пустой макет
        if background_image:
            # Добавление изображения в качестве фона
            slide.shapes.add_picture(background_image, 0, 0, width=prs.slide_width, height=prs.slide_height)

        # Добавление заголовка и текста
        title = slide.shapes.title
        title.text = slide_data["title"]
        content = slide.shapes.placeholders[1].text_frame
        content.text = slide_data["content"]

        # Добавление изображения, если оно есть
        if slide_data["image"]:
            img_path = slide_data["image"]
            slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(4))

    # Сохраняем презентацию
    prs.save(f"{topic}_presentation.pptx")
    return f"{topic}_presentation.pptx"


# конвертация PPTX в PDF
def convert_pptx_to_pdf(pptx_path):
    pdf_path = pptx_path.replace(".pptx", ".pdf")
    # TODO - конвертация на любой ОС
    os.system(f"libreoffice --headless --convert-to pdf {pptx_path} --outdir {os.path.dirname(pptx_path)}")
    return pdf_path


# Создание миниатюр слайдов
def create_slide_previews(pptx_path):
    images = convert_from_path(pptx_path)
    return images


